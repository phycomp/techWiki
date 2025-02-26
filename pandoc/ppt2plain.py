#!/usr/bin/env python
from streamlit import sidebar, radio as stRadio, text_input, tabs as stTAB, markdown as stMarkdown
#import io
from pathlib import Path
from os.path import dirname, basename
from pptx2txt2 import extract_images as PPTxtrctIMG, extract_text as xtrctPPT, extract_text_per_slide as xtrctSlide
from docx2txt2 import extract_text as xtrctDoc, extract_images as xtrctDocIMG
from pptx import Presentation
from stUtil import rndrCode
from fhndlrUtil import rcrsvDIR, profileDIR
from mimetypes import guess_type
from base64 import b64encode

MENU, 表單=[], ['轉換PPT', '轉換ODP']	#, '錯綜複雜', '二十四節氣'
for ndx, Menu in enumerate(表單): MENU.append(f'{ndx}{Menu}')
with sidebar:
  menu=stRadio('表單', MENU, horizontal=True, index=0)
  來源=text_input('來源')
  目的=text_input('目的')
  #menu=stRadio('表單', MENU, horizontal=True, index=0)
brwsrType=['PPT', 'DOC', 'PDF']
tabPPT, tabDOC, tabPDF = stTAB(brwsrType)
母層=Path('/data/GITs/md2pdf')
根層=母層.iterdir()
temp, pptx, docx=[], None, None
with sidebar:
  for x in 根層:
    if x.is_dir():temp.append(x)
  neatDIR=map(basename, temp)
  目錄=stRadio('目錄', neatDIR, horizontal=True, index=0)

def parsePRFL(dirPRFL):
  tmp={}
  for k, v in dirPRFL.items(): tmp.update({k:len(v)})
  return tmp

if menu==len(表單):
  pass
elif menu==MENU[1]:
  tblName='sutra'
elif menu==MENU[0]:
  if 目錄:
    目錄=母層/目錄  #f'{母層}/{目錄}'
    for 根, 檔案  in rcrsvDIR(DIR=目錄):
      根=Path(根)
      with sidebar:
        dirPRFL=profileDIR(根, 檔案)
        rndrCode(parsePRFL(dirPRFL))
    with tabPPT:
      rndrCode(['powerpoint' in dirPRFL])
      分頁檔案=[]
      for k in dirPRFL:
        if k is not None and ('presentation' in k or 'powerpoint' in k):
      #if 分頁檔案:=dirPRFL.get('application/vnd.ms-powerpoint'):    #zip
          #分頁檔案=dirPRFL.get(k)#'application/vnd.ms-powerpoint')
          分頁檔案+=dirPRFL.get(k)#('application/vnd.openxmlformats-officedocument.presentationml.presentation')
      根=Path(dirname(分頁檔案[0]))
      with sidebar:
        neatPPT=map(basename, 分頁檔案)
        prsFILE=stRadio('PPTX', neatPPT, horizontal=True, index=None)
      if prsFILE:
        fullPPT=根/prsFILE #.as_posix()
        #ppt = Presentation()
        #  from gzinfo import read_gz_info
        #  gzINFO = read_gz_info(fullPPT)
        #  rndrCode(gzINFO.fname)
        #  from gzip import GZipFile
        #  gz=GZipFile(fullPPT)
        #except Exception as e:
        try:
          from io import BytesIO
          with open(fullPPT.resolve(), 'rb') as fin:
            prsSRC = BytesIO(fin.read())
            prs = Presentation(prsSRC) #.resolve(strict=True)
          pptINFO=[]
          for slide in prs.slides:
            for shape in slide.shapes:
              if hasattr(shape, "text"):
                pptINFO.append(shape.text)
          stMarkdown(pptINFO, unsafe_allow_html=False) #.split('\n')'\n'.join(pptINFO)
        except:
          pass
          #from pptdump import PPTDumper
          #dumper = PPTDumper(fullPPT)   #, globals.params
          #rsltINFO=dumper.dump()
          #rndrCode(rsltINFO)    #dumptext()
        finally:
          cntxt = xtrctPPT(fullPPT) #'/data/GITs/md2pdf/newDOCs/六祖壇經三.pptx''/data/GITs/md2pdf/歷程/SurvivalAnalysis.ppt'
          stMarkdown(cntxt, unsafe_allow_html=False) #.split('\n')
          #cntxtSlide = xtrctSlide(fullPPT)
          #stMarkdown(cntxtSlide, unsafe_allow_html=False)
          #imgPath = PPTxtrctIMG(fullPPT, 'outIMG')
          #from odf.presentation import Notes
          #if not dumper.dump(): print("FAILURE\n")
          #if globals.params.dumpText:
          #    globals.dumptext()

          #from odf.opendocument import load
          #from odf import text
          #doc=load(fullPPT)
          #slides=doc.presentation
          #notes=slides.getElementsByType(Notes)
          #for page in notes:
          #   pars = page.getElementsByType(text.P)
          #   for par in pars:
          #       rndrCode(['par', par])

          #from spire.presentation import *
          #from spire.presentation.common import * #FileFormat
                #file.write(shape.text+"\n")
        #iter_picture_shapes(prs)
        #rndrCode(fullPPT)
    with tabDOC:
      #from odf.opendocument import load as odfLoad
      #from odf import text as odfTXT, teletype
      #from odf.opendocument import load as odfLoad
      if 分頁檔案:=dirPRFL.get('application/msword'):
        根=Path(dirname(分頁檔案[0]))
        #from spire.doc import Document# import LoadFromFile .common
        with sidebar:
          neatDOC=map(basename, 分頁檔案)
          docx=stRadio('DOCX', neatDOC, horizontal=True, index=None)
      if docx:
        fullDOC=根/docx
        doc = Document()
        doc.LoadFromFile(fullDOC.as_posix())
        docCntxt = doc.GetText()
        rndrCode(docCntxt)
        #doc.close()
        #docTxt = xtrctDoc(fullPath.as_posix())
        #docIMGpath= xtrctDocIMG(fullPath, 'outIMG')
        #rndrCode([docTxt, docIMGpath])
        #from spire.doc import *

# Create a Document object
# Load a Word document

# Extract the text of the document

# Write the extracted text into a text file
        #with open("Output/DocumentText.txt", "w", encoding="utf-8") as file:
        #     file.write(document_text)

        #textdoc = odfLoad(fullPath)  #"your.odt"
        #allParas = textdoc.getElementsByType(odfTXT.P)
        #out=teletype.extractText(allParas[0])
        #with open(f'{base}.raw', 'w') as fout:
        #  fout.write(out)
        #rndrCode(fullPath)
        #cntxt = xtrctTXT(fullPath.as_posix())
        #rndrCode(cntxt)
        #cntxtSlide = xtrctSlide(pptx)
        #rndrCode(cntxtSlide)
        #imgPath = xtrctIMG(pptx)
      #print('\n'.join(cntxt.split('\n')))    #imgPath, cntxtSlide, '\n', 

# actual Paths
#pptx_path = Path(__file__).parent / "my.pptx"
#imageOut = Path(__file__).parent / "my" / "images"
#imageOut.mkdir(parents=True)

#text2 = xtrctTXT(pptx_path)
#text_per_slide2 = xtrctSlide(pptx_path)
#image_paths2 = xtrctIMG(pptx_path, image_out)

# bytestreams
#pptx_bytes = b"..."
#bytes_io = io.BytesIO(pptx_bytes)
#text3 = pptx2txt2.extract_text(bytes_io)
#text_per_slide3 = pptx2txt2.extract_text_per_slide(bytes_io)
#image_paths3 = pptx2txt2.extract_images(bytes_io, "path/to/images/out")
